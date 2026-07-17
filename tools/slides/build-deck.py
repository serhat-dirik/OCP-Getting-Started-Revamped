#!/usr/bin/env python3
"""Build the instructor PPTX deck from slides/outlines/*.md (ADR-0005).

Outline schema (04-STYLE-GUIDE §6):
    # <deck/module title>          -> section title slide
    ## Slide: <title>              -> one slide
    - bullet (<=5 per slide)       -> content placeholder
    Notes: <paragraph>             -> speaker notes
    Visual: <diagram spec/ref>     -> speaker-note annotation [VISUAL: ...]

Template swap (gate answer 2026-07-08: placeholder until the real Red Hat file lands):
    slides/template/reference.pptx|.potx  -> used when present (Red Hat template drop-in)
    otherwise                             -> python-pptx built-in default template
Layout names are resolved via slides/template/layout-map.yaml when present
(maps role -> layout name), else first-match against common names.

Usage:
    build-deck.py --out slides/dist/deck.pptx [outline.md ...]
    (no outline args = every slides/outlines/*.md, sorted)
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:  # pragma: no cover
    sys.exit("❌ python-pptx not installed — run: pip install python-pptx (CI does this)")

REPO_ROOT = Path(__file__).resolve().parents[2]
TEMPLATE_DIR = REPO_ROOT / "slides" / "template"

# Role -> candidate layout names (first present wins). Overridable via layout-map.yaml.
LAYOUT_CANDIDATES = {
    "section": ["Section Header", "Section Title", "Title Slide"],
    "content": ["Title and Content", "Title, Content", "Content with Caption"],
}


def find_template() -> Path | None:
    if TEMPLATE_DIR.is_dir():
        for pattern in ("*.potx", "*.pptx"):
            hits = sorted(TEMPLATE_DIR.glob(pattern))
            if hits:
                return hits[0]
    return None


def load_layout_map() -> dict[str, str]:
    mapping_file = TEMPLATE_DIR / "layout-map.yaml"
    mapping: dict[str, str] = {}
    if mapping_file.is_file():
        for line in mapping_file.read_text().splitlines():
            line = line.strip()
            if line and not line.startswith("#") and ":" in line:
                role, name = line.split(":", 1)
                mapping[role.strip()] = name.strip().strip("'\"")
    return mapping


def pick_layout(prs: Presentation, role: str, layout_map: dict[str, str]):
    names = {layout.name: layout for layout in prs.slide_layouts}
    if role in layout_map and layout_map[role] in names:
        return names[layout_map[role]]
    for candidate in LAYOUT_CANDIDATES[role]:
        if candidate in names:
            return names[candidate]
    return prs.slide_layouts[1 if role == "content" else 0]  # positional fallback


def parse_outline(path: Path) -> list[dict]:
    """Return slide dicts: {role, title, bullets, notes, visual}."""
    slides: list[dict] = []
    current: dict | None = None
    for raw in path.read_text().splitlines():
        line = raw.rstrip()
        if m := re.match(r"^#\s+(?!#)(.+)$", line):
            slides.append({"role": "section", "title": m.group(1).strip(),
                           "bullets": [], "notes": "", "visual": ""})
            current = slides[-1]
        elif m := re.match(r"^##\s*Slide:\s*(.+)$", line):
            slides.append({"role": "content", "title": m.group(1).strip(),
                           "bullets": [], "notes": "", "visual": ""})
            current = slides[-1]
        elif current and (m := re.match(r"^-\s+(.+)$", line)):
            current["bullets"].append(m.group(1).strip())
        elif current and (m := re.match(r"^Notes:\s*(.*)$", line)):
            current["notes"] = m.group(1).strip()
        elif current and (m := re.match(r"^Visual:\s*(.*)$", line)):
            current["visual"] = m.group(1).strip()
        elif current and current["notes"] and line and not line.startswith(("#", "-")):
            current["notes"] += " " + line.strip()  # continuation lines
    return slides


def add_slide(prs: Presentation, spec: dict, layout_map: dict[str, str]) -> None:
    slide = prs.slides.add_slide(pick_layout(prs, spec["role"], layout_map))
    if slide.shapes.title is not None:
        slide.shapes.title.text = spec["title"]
    body = next((ph for ph in slide.placeholders
                 if ph.placeholder_format.idx != 0 and ph.has_text_frame), None)
    if body is not None and spec["bullets"]:
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(spec["bullets"][:5]):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = bullet
            para.font.size = Pt(20)
    notes = spec["notes"] + (f"\n[VISUAL: {spec['visual']}]" if spec["visual"] else "")
    if notes.strip():
        slide.notes_slide.notes_text_frame.text = notes.strip()


def ordered_outlines() -> list[Path]:
    """Deck order comes from /modules.yaml (position = number, decoupled from slug): one
    slides/outlines/<slug>.md per module, in catalog order. Falls back to alphabetical when the
    manifest is absent; warns (but still includes) outlines missing from / not listed in it."""
    outdir = REPO_ROOT / "slides" / "outlines"
    manifest = REPO_ROOT / "modules.yaml"
    all_md = sorted(outdir.glob("*.md"))
    if not manifest.is_file():
        print("⚠ modules.yaml not found — falling back to alphabetical outline order")
        return all_md
    slugs = re.findall(r"^\s*-\s*slug:\s*(\S+)", manifest.read_text(), re.M)
    ordered: list[Path] = []
    seen: set[str] = set()
    for slug in slugs:
        p = outdir / f"{slug}.md"
        if p.is_file():
            ordered.append(p)
            seen.add(p.name)
        else:
            print(f"⚠ no outline for module '{slug}' ({p.name}) — skipping")
    for p in all_md:  # keep any stray outline not in the manifest, at the end
        if p.name not in seen:
            print(f"⚠ outline {p.name} not listed in modules.yaml — appended at end")
            ordered.append(p)
    return ordered


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("outlines", nargs="*", type=Path)
    ap.add_argument("--out", type=Path, default=REPO_ROOT / "slides" / "dist" / "deck.pptx")
    args = ap.parse_args()

    outlines = args.outlines or ordered_outlines()
    if not outlines:
        print("❌ no outlines found under slides/outlines/")
        return 1

    template = find_template()
    prs = Presentation(str(template)) if template else Presentation()
    layout_map = load_layout_map()
    print(f"▶ template: {template.name if template else 'python-pptx default (placeholder mode)'}")

    total = 0
    for outline in outlines:
        for spec in parse_outline(outline):
            add_slide(prs, spec, layout_map)
            total += 1
        print(f"  ✓ {outline.name}")

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))

    # self-check: reopen and assert the deck is intact (CI gate per ADR-0005)
    reopened = Presentation(str(args.out))
    built = len(reopened.slides.__iter__.__self__._sldIdLst)  # stable slide count
    if built < total:
        print(f"❌ self-check failed: built {built} < parsed {total}")
        return 1
    print(f"✅ {built} slides -> {args.out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
